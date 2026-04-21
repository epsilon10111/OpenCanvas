#!/usr/bin/env python3
"""
课程知识库 - 读取课程文件，提取文本内容作为 AI 参考

支持格式:
- .pdf (PyPDF2)
- .txt, .md
- .docx (python-docx)
- .pptx (python-pptx)
"""

import os
from pathlib import Path
from typing import Any

KB_MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB per file
KB_MAX_TOTAL_TEXT = 200_000  # 最大提取字符数


def _try_import(name: str) -> bool:
    try:
        __import__(name)
        return True
    except ImportError:
        return False


def _read_pdf(path: Path, max_chars: int = 10000) -> str:
    """读取 PDF 文本内容"""
    if _try_import("PyPDF2"):
        import PyPDF2
        text_parts = []
        try:
            with open(path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages[:50]:  # 最多50页
                    t = page.extract_text() or ""
                    text_parts.append(t)
                    if len("".join(text_parts)) > max_chars:
                        break
        except Exception as e:
            return f"[PDF读取失败: {e}]"
        return "\n\n".join(text_parts)[:max_chars]
    return "[PDF: 请安装 PyPDF2]"


def _read_docx(path: Path, max_chars: int = 10000) -> str:
    """读取 DOCX 文本内容"""
    if _try_import("docx"):
        import docx
        try:
            doc = docx.Document(str(path))
            text = "\n".join(p.text for p in doc.paragraphs)
            return text[:max_chars]
        except Exception as e:
            return f"[DOCX读取失败: {e}]"
    return "[DOCX: 请安装 python-docx]"


def _read_pptx(path: Path, max_chars: int = 10000) -> str:
    """读取 PPTX 文本内容"""
    if _try_import("pptx"):
        from pptx import Presentation
        try:
            prs = Presentation(str(path))
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text_parts.append(para.text)
            return "\n".join(text_parts)[:max_chars]
        except Exception as e:
            return f"[PPTX读取失败: {e}]"
    return "[PPTX: 请安装 python-pptx]"


def _read_text_file(path: Path, max_chars: int = 10000) -> str:
    """读取纯文本/Markdown"""
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()[:max_chars]
    except Exception as e:
        return f"[文本读取失败: {e}]"


def read_course_file(path: Path) -> tuple[str, str]:
    """
    读取单个课程文件
    返回: (文件路径, 提取的文本内容)
    """
    if not path.exists() or path.is_dir():
        return str(path), ""

    if path.stat().st_size > KB_MAX_FILE_SIZE:
        return str(path), f"[文件过大，跳过: {path.stat().st_size / 1024 / 1024:.1f}MB]"

    suffix = path.suffix.lower()

    if suffix == ".pdf":
        return str(path), _read_pdf(path)
    elif suffix in (".txt", ".md", ".py", ".java", ".c", ".cpp", ".js", ".ts"):
        return str(path), _read_text_file(path)
    elif suffix == ".docx":
        return str(path), _read_docx(path)
    elif suffix == ".pptx":
        return str(path), _read_pptx(path)
    else:
        return str(path), ""


def build_course_knowledge(
    course_dir: Path,
    max_files: int = 20,
    max_chars_per_file: int = 8000,
) -> list[dict[str, str]]:
    """
    构建课程知识库
    扫描课程目录，提取文本内容
    返回: [{"file": str, "content": str}, ...]
    """
    if not course_dir.exists():
        return []

    knowledge = []
    supported_exts = {".pdf", ".txt", ".md", ".py", ".java", ".c", ".cpp",
                      ".js", ".ts", ".docx", ".pptx"}

    # 遍历课程目录
    for root, dirs, files in os.walk(course_dir):
        # 跳过隐藏目录和 __pycache__
        dirs[:] = [d for d in dirs if not d.startswith(".") and d != "__pycache__"]
        for fname in sorted(files):
            if len(knowledge) >= max_files:
                break
            fpath = Path(root) / fname
            if fpath.suffix.lower() in supported_exts:
                rel_path = fpath.relative_to(course_dir.parent)
                fp, content = read_course_file(fpath)
                if content and not content.startswith("["):
                    # 截断到指定长度
                    content = content[:max_chars_per_file]
                    knowledge.append({
                        "file": str(rel_path),
                        "content": content,
                    })
            if len(knowledge) >= max_files:
                break

    return knowledge


def check_course_files(course_dir: Path) -> dict[str, Any]:
    """
    检查课程文件状态
    返回: {"exists": bool, "file_count": int, "total_size": int, "has_content": bool}
    """
    if not course_dir.exists():
        return {"exists": False, "file_count": 0, "total_size": 0, "has_content": False}

    file_count = 0
    total_size = 0
    has_content = False

    for root, dirs, files in os.walk(course_dir):
        dirs[:] = [d for d in dirs if not d.startswith(".") and d != "__pycache__"]
        for fname in files:
            fpath = Path(root) / fname
            file_count += 1
            total_size += fpath.stat().st_size
            if not has_content and fpath.suffix.lower() in {
                ".pdf", ".txt", ".md", ".docx", ".pptx"
            }:
                has_content = True

    return {
        "exists": True,
        "file_count": file_count,
        "total_size": total_size,
        "has_content": has_content,
    }


def ensure_course_files(
    course_id: int,
    base_url: str,
    token: str,
    course_dir: Path,
) -> dict[str, Any]:
    """
    确保课程文件已下载
    如果目录为空，尝试从 Canvas API 下载
    返回: 文件状态信息
    """
    # 先检查
    status = check_course_files(course_dir)
    if status["has_content"] and status["file_count"] > 0:
        return status

    # 文件不足，尝试下载
    try:
        from download_courses import (
            canvas_auth_headers,
            download_course_files,
            sanitize_path_segment,
        )
        import httpx

        headers = canvas_auth_headers(token)
        course_dir.mkdir(parents=True, exist_ok=True)

        with httpx.Client(timeout=120.0, follow_redirects=True) as client:
            ok, bad, over = download_course_files(
                client, base_url, headers, course_id, course_dir
            )
            status["downloaded"] = ok
            status["failed"] = bad
            status["over_limit"] = over
    except ImportError:
        status["download_error"] = "无法导入下载模块"
    except Exception as e:
        status["download_error"] = str(e)

    return status
